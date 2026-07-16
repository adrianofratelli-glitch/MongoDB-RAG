import os
import sys
import json
import time
import argparse
from pathlib import Path
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
import voyageai
from config import DB_NAME, CLIENT_ID
from db import get_client
from dotenv import load_dotenv

load_dotenv()

SUPPORTED_FORMATS = {
    ".pdf", ".docx", ".doc", ".txt", ".csv",
    ".md", ".markdown", ".html", ".htm", ".json", ".xlsx", ".xls", ".pptx",
}


class JSONTextLoader:
    """One Document per top-level record (or the whole file if it's a single object)."""

    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        data = json.loads(Path(self.file_path).read_text(encoding="utf-8"))
        records = data if isinstance(data, list) else [data]
        return [
            Document(
                page_content=json.dumps(r, ensure_ascii=False, indent=2),
                metadata={"page": i},
            )
            for i, r in enumerate(records)
        ]


class ExcelTextLoader:
    """One Document per worksheet, rendered as tab-separated rows."""

    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        from openpyxl import load_workbook
        wb = load_workbook(self.file_path, data_only=True, read_only=True)
        docs = []
        for i, sheet in enumerate(wb.worksheets):
            lines = [
                "\t".join("" if c is None else str(c) for c in row)
                for row in sheet.iter_rows(values_only=True)
            ]
            docs.append(Document(
                page_content="\n".join(lines),
                metadata={"page": i, "sheet": sheet.title},
            ))
        return docs


class PowerPointTextLoader:
    """One Document per slide."""

    def __init__(self, file_path):
        self.file_path = file_path

    def load(self):
        from pptx import Presentation
        prs = Presentation(self.file_path)
        docs = []
        for i, slide in enumerate(prs.slides):
            texts = [
                shape.text for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text
            ]
            docs.append(Document(page_content="\n".join(texts), metadata={"page": i}))
        return docs


def get_loader(file_path: str):
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        from langchain_community.document_loaders import PyPDFLoader
        return PyPDFLoader(file_path)
    if ext in (".docx", ".doc"):
        from langchain_community.document_loaders import Docx2txtLoader
        return Docx2txtLoader(file_path)
    if ext in (".txt", ".md", ".markdown"):
        from langchain_community.document_loaders import TextLoader
        return TextLoader(file_path, encoding="utf-8")
    if ext == ".csv":
        from langchain_community.document_loaders import CSVLoader
        return CSVLoader(file_path)
    if ext in (".html", ".htm"):
        from langchain_community.document_loaders import BSHTMLLoader
        return BSHTMLLoader(file_path, open_encoding="utf-8")
    if ext == ".json":
        return JSONTextLoader(file_path)
    if ext in (".xlsx", ".xls"):
        return ExcelTextLoader(file_path)
    if ext == ".pptx":
        return PowerPointTextLoader(file_path)
    raise ValueError(
        f"Unsupported format: '{ext}'. "
        f"Accepted formats: {', '.join(sorted(SUPPORTED_FORMATS))}"
    )


def ingest(file_path: str, reset: bool = False, nivel_acesso: str = "publico") -> None:
    path = Path(file_path)
    if not path.exists():
        print(f"File not found: {file_path}")
        sys.exit(1)

    client = get_client()
    collection = client[DB_NAME]["documents"]

    source_name = path.stem
    existing = collection.count_documents({"metadata.source": source_name})

    if existing > 0 and not reset:
        print(
            f"'{source_name}' is already indexed ({existing} chunks). "
            "Use --reset to re-index."
        )
        return

    if reset and existing > 0:
        print(f"Removing {existing} chunks from '{source_name}'...")
        collection.delete_many({"metadata.source": source_name})

    voyage = voyageai.Client(api_key=os.environ["VOYAGE_API_KEY"])

    print(f"Loading {path.name}...")
    loader = get_loader(file_path)
    docs = loader.load()
    print(f"   {len(docs)} pages/sections loaded")

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=150,
        separators=["\n\n", "\n", ".", " "],
    )
    chunks = splitter.split_documents(docs)
    print(f"   {len(chunks)} chunks generated")

    texts = [c.page_content for c in chunks]
    batch_size = 10  # conservative for the free tier (10K TPM)
    docs_to_insert = []

    print("Generating embeddings (free tier may take several minutes for large documents)...")
    for i in range(0, len(texts), batch_size):
        batch_texts = texts[i : i + batch_size]
        batch_chunks = chunks[i : i + batch_size]

        result = voyage.embed(batch_texts, model="voyage-3", input_type="document")

        for j, embedding in enumerate(result.embeddings):
            docs_to_insert.append({
                "text": batch_chunks[j].page_content,
                "embedding": embedding,
                "metadata": {
                    "source": source_name,
                    "client_id": CLIENT_ID,
                    "file": path.name,
                    "page": batch_chunks[j].metadata.get("page", 0),
                    "chunk_id": i + j,
                    # Access level used by the $vectorSearch / Atlas Search filter.
                    "nivel_acesso": nivel_acesso,
                },
            })

        progress = min(i + batch_size, len(texts))
        print(f"   {progress}/{len(texts)} chunks | batch {i // batch_size + 1}", end="\r")

        # Rate limit: 3 RPM, i.e. one request every 20s (with margin)
        if i + batch_size < len(texts):
            time.sleep(22)

    print(f"\nInserting {len(docs_to_insert)} documents into Atlas (DB: {DB_NAME})...")
    collection.insert_many(docs_to_insert)
    print("Ingestion complete.")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Ingest documents into the RAG store "
        "(PDF, DOCX, TXT, CSV, Markdown, HTML, JSON, XLSX, PPTX)"
    )
    parser.add_argument(
        "file",
        nargs="?",
        default="data/documento.pdf",
        help="Path to the file to index",
    )
    parser.add_argument(
        "--reset",
        action="store_true",
        help="Remove existing chunks and re-index the document",
    )
    parser.add_argument(
        "--nivel",
        choices=["publico", "restrito"],
        default="publico",
        help="Access level (public/restricted) assigned to this document's chunks",
    )
    args = parser.parse_args()
    ingest(args.file, reset=args.reset, nivel_acesso=args.nivel)
